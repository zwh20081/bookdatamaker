"""Document parser for PDF, EPUB, and PPTX files."""

import contextlib
import os
import sys
from pathlib import Path
from typing import Iterator, List, Tuple

from PIL import Image


@contextlib.contextmanager
def _suppress_mupdf_stderr():
    """Suppress MuPDF C-level stderr warnings (e.g. 'No common ancestor in structure tree')."""
    try:
        devnull_fd = os.open(os.devnull, os.O_WRONLY)
        old_stderr_fd = os.dup(2)
        os.dup2(devnull_fd, 2)
        os.close(devnull_fd)
        yield
    except OSError:
        yield  # If redirection fails, just proceed normally
    else:
        os.dup2(old_stderr_fd, 2)
        os.close(old_stderr_fd)


class DocumentParser:
    """Parse PDF, EPUB, and PPTX documents into images or text."""

    @staticmethod
    def is_pdf(file_path: Path) -> bool:
        """Check if file is PDF."""
        return file_path.suffix.lower() == ".pdf"

    @staticmethod
    def is_epub(file_path: Path) -> bool:
        """Check if file is EPUB."""
        return file_path.suffix.lower() == ".epub"

    @staticmethod
    def is_pptx(file_path: Path) -> bool:
        """Check if file is PPTX."""
        return file_path.suffix.lower() == ".pptx"

    @staticmethod
    def parse_pdf_to_images(pdf_path: Path, dpi: int = 200) -> List[Tuple[int, Image.Image]]:
        """Convert PDF pages to images.

        Args:
            pdf_path: Path to PDF file
            dpi: Resolution for rendering (default: 200)

        Returns:
            List of tuples (page_number, image)
        """
        try:
            import fitz  # PyMuPDF

            with _suppress_mupdf_stderr():
                doc = fitz.open(pdf_path)
            images = []

            for page_num in range(len(doc)):
                page = doc[page_num]
                # Render page to pixmap
                with _suppress_mupdf_stderr():
                    pix = page.get_pixmap(dpi=dpi)
                # Convert to PIL Image
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                images.append((page_num + 1, img))

            doc.close()
            return images

        except ImportError:
            raise ImportError(
                "PyMuPDF (fitz) is required for PDF parsing. "
                "Install with: pip install pymupdf"
            )

    @staticmethod
    def iter_pdf_images(pdf_path: Path, dpi: int = 200, start_page: int = 1) -> Iterator[Tuple[int, Image.Image]]:
        """Yield PDF pages as images one at a time.

        Args:
            pdf_path: Path to PDF file
            dpi: Resolution for rendering
            start_page: 1-based page to start from

        Yields:
            Tuples of (page_number, image)
        """
        try:
            import fitz  # PyMuPDF

            with _suppress_mupdf_stderr():
                doc = fitz.open(pdf_path)
            try:
                for page_index in range(max(start_page - 1, 0), len(doc)):
                    page = doc[page_index]
                    with _suppress_mupdf_stderr():
                        pix = page.get_pixmap(dpi=dpi)
                    image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    yield page_index + 1, image
            finally:
                doc.close()

        except ImportError:
            raise ImportError(
                "PyMuPDF (fitz) is required for PDF parsing. "
                "Install with: pip install pymupdf"
            )

    @staticmethod
    def parse_pdf_text(pdf_path: Path) -> List[Tuple[int, str]]:
        """Extract text directly from PDF (if available).

        Args:
            pdf_path: Path to PDF file

        Returns:
            List of tuples (page_number, text)
        """
        try:
            import fitz  # PyMuPDF

            with _suppress_mupdf_stderr():
                doc = fitz.open(pdf_path)
            pages_text = []

            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                pages_text.append((page_num + 1, text))

            doc.close()
            return pages_text

        except ImportError:
            raise ImportError(
                "PyMuPDF (fitz) is required for PDF parsing. "
                "Install with: pip install pymupdf"
            )
    
    @staticmethod
    def parse_pdf_text_and_images(pdf_path: Path, dpi: int = 200) -> List[Tuple[int, tuple[str, Image.Image]]]:
        """Extract both text and images from PDF.

        Args:
            pdf_path: Path to PDF file
            dpi: DPI for rendering images

        Returns:
            List of tuples (page_number, (text, image))
        """
        try:
            import fitz  # PyMuPDF

            with _suppress_mupdf_stderr():
                doc = fitz.open(pdf_path)
            pages = []

            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Extract text
                text = page.get_text()
                
                # Render page as image
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                with _suppress_mupdf_stderr():
                    pix = page.get_pixmap(matrix=mat)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                pages.append((page_num + 1, (text, img)))

            doc.close()
            return pages

        except ImportError:
            raise ImportError(
                "PyMuPDF (fitz) is required for PDF parsing. "
                "Install with: pip install pymupdf"
            )

    @staticmethod
    def iter_pdf_text_and_images(
        pdf_path: Path, dpi: int = 200, start_page: int = 1
    ) -> Iterator[Tuple[int, tuple[str, Image.Image]]]:
        """Yield PDF pages as (text, image) tuples one at a time.

        Args:
            pdf_path: Path to PDF file
            dpi: DPI for rendering images
            start_page: 1-based page to start from

        Yields:
            Tuples of (page_number, (text, image))
        """
        try:
            import fitz  # PyMuPDF

            with _suppress_mupdf_stderr():
                doc = fitz.open(pdf_path)
            try:
                for page_index in range(max(start_page - 1, 0), len(doc)):
                    page = doc[page_index]
                    text = page.get_text()
                    mat = fitz.Matrix(dpi / 72, dpi / 72)
                    with _suppress_mupdf_stderr():
                        pix = page.get_pixmap(matrix=mat)
                    image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    yield page_index + 1, (text, image)
            finally:
                doc.close()

        except ImportError:
            raise ImportError(
                "PyMuPDF (fitz) is required for PDF parsing. "
                "Install with: pip install pymupdf"
            )

    @staticmethod
    def get_pdf_page_count(pdf_path: Path) -> int:
        """Get the number of pages in a PDF file."""
        try:
            import fitz  # PyMuPDF

            with _suppress_mupdf_stderr():
                doc = fitz.open(pdf_path)
            try:
                return len(doc)
            finally:
                doc.close()

        except ImportError:
            raise ImportError(
                "PyMuPDF (fitz) is required for PDF parsing. "
                "Install with: pip install pymupdf"
            )

    @staticmethod
    def parse_epub_to_images(epub_path: Path, width: int = 800) -> List[Tuple[int, Image.Image]]:
        """Convert EPUB pages to images.

        Args:
            epub_path: Path to EPUB file
            width: Image width in pixels

        Returns:
            List of tuples (page_number, image)
        """
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            from PIL import ImageDraw, ImageFont

            book = epub.read_epub(epub_path)
            images = []
            page_num = 1

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    # Parse HTML content
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text = soup.get_text()

                    # Create simple text image (for OCR simulation)
                    # In production, you might want to render HTML properly
                    img = Image.new("RGB", (width, 1000), color="white")
                    draw = ImageDraw.Draw(img)

                    try:
                        font = ImageFont.truetype("arial.ttf", 16)
                    except:
                        font = ImageFont.load_default()

                    # Simple text rendering
                    y_offset = 10
                    for line in text.split("\n")[:50]:  # Limit lines
                        if line.strip():
                            draw.text((10, y_offset), line[:100], fill="black", font=font)
                            y_offset += 20

                    images.append((page_num, img))
                    page_num += 1

            return images

        except ImportError:
            raise ImportError(
                "ebooklib and beautifulsoup4 are required for EPUB parsing. "
                "Install with: pip install ebooklib beautifulsoup4"
            )

    @staticmethod
    def parse_epub_text(epub_path: Path) -> List[Tuple[int, str]]:
        """Extract text directly from EPUB.

        Args:
            epub_path: Path to EPUB file

        Returns:
            List of tuples (chapter_number, text)
        """
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(epub_path)
            pages_text = []
            page_num = 1

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text = soup.get_text()
                    pages_text.append((page_num, text))
                    page_num += 1

            return pages_text

        except ImportError:
            raise ImportError(
                "ebooklib and beautifulsoup4 are required for EPUB parsing. "
                "Install with: pip install ebooklib beautifulsoup4"
            )
    
    @staticmethod
    def parse_epub_text_and_images(epub_path: Path, width: int = 800) -> List[Tuple[int, tuple[str, Image.Image]]]:
        """Extract both text and rendered images from EPUB.

        Args:
            epub_path: Path to EPUB file
            width: Image width in pixels

        Returns:
            List of tuples (chapter_number, (text, image))
        """
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            from PIL import ImageDraw, ImageFont

            book = epub.read_epub(epub_path)
            pages = []
            page_num = 1

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    # Extract text
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text = soup.get_text()

                    # Create simple text image for consistency
                    img = Image.new("RGB", (width, 1000), color="white")
                    draw = ImageDraw.Draw(img)

                    try:
                        font = ImageFont.truetype("arial.ttf", 16)
                    except:
                        font = ImageFont.load_default()

                    # Draw text on image (simplified rendering)
                    y_offset = 20
                    for line in text[:2000].split("\n"):  # Limit text for rendering
                        if y_offset > 950:
                            break
                        draw.text((20, y_offset), line[:80], fill="black", font=font)
                        y_offset += 20

                    pages.append((page_num, (text, img)))
                    page_num += 1

            return pages

        except ImportError:
            raise ImportError(
                "ebooklib and beautifulsoup4 are required for EPUB parsing. "
                "Install with: pip install ebooklib beautifulsoup4"
            )

    @staticmethod
    def get_epub_page_count(epub_path: Path) -> int:
        """Get the number of document items in an EPUB file."""
        try:
            import ebooklib
            from ebooklib import epub

            book = epub.read_epub(epub_path)
            return sum(1 for item in book.get_items() if item.get_type() == ebooklib.ITEM_DOCUMENT)

        except ImportError:
            raise ImportError(
                "ebooklib and beautifulsoup4 are required for EPUB parsing. "
                "Install with: pip install ebooklib beautifulsoup4"
            )

    @staticmethod
    def parse_pptx_to_images(
        pptx_path: Path, width: int = 1280
    ) -> List[Tuple[int, Image.Image]]:
        """Render PPTX slides as text images.

        Args:
            pptx_path: Path to PPTX file
            width: Image width in pixels

        Returns:
            List of tuples (slide_number, image)
        """
        try:
            from pptx import Presentation
            from pptx.util import Emu
            from PIL import ImageDraw, ImageFont

            prs = Presentation(str(pptx_path))
            slide_width = prs.slide_width or Emu(9144000)  # default 10 inches
            slide_height = prs.slide_height or Emu(6858000)  # default 7.5 inches
            aspect = int(slide_height) / int(slide_width)
            height = int(width * aspect)

            pages: List[Tuple[int, Image.Image]] = []

            for idx, slide in enumerate(prs.slides, start=1):
                text_parts: List[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                text_parts.append(line)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip(" |" ):
                                text_parts.append(row_text)

                img = Image.new("RGB", (width, height), color="white")
                draw = ImageDraw.Draw(img)
                try:
                    font = ImageFont.truetype("arial.ttf", 18)
                except Exception:
                    font = ImageFont.load_default()

                y_offset = 20
                for line in text_parts:
                    if y_offset > height - 30:
                        break
                    draw.text((20, y_offset), line[:120], fill="black", font=font)
                    y_offset += 24

                pages.append((idx, img))

            return pages

        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX parsing. "
                "Install with: pip install python-pptx"
            )

    @staticmethod
    def parse_pptx_text(pptx_path: Path) -> List[Tuple[int, str]]:
        """Extract text from PPTX slides.

        Args:
            pptx_path: Path to PPTX file

        Returns:
            List of tuples (slide_number, text)
        """
        try:
            from pptx import Presentation
            import io

            prs = Presentation(str(pptx_path))
            pages: List[Tuple[int, str]] = []

            for idx, slide in enumerate(prs.slides, start=1):
                text_parts: List[str] = []
                img_idx = 0
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                text_parts.append(line)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip(" |"):
                                text_parts.append(row_text)
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            shape.image  # noqa: verify image exists
                            text_parts.append(f"![](images/{img_idx}.jpg)")
                            img_idx += 1
                        except Exception:
                            pass

                pages.append((idx, "\n".join(text_parts)))

            return pages

        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX parsing. "
                "Install with: pip install python-pptx"
            )

    @staticmethod
    def parse_pptx_text_and_images(
        pptx_path: Path, width: int = 1280
    ) -> List[Tuple[int, tuple]]:
        """Extract both text and rendered images from PPTX.

        Args:
            pptx_path: Path to PPTX file
            width: Image width in pixels

        Returns:
            List of tuples (slide_number, (text, rendered_image, embedded_images))
            where embedded_images is a list of PIL Images extracted from the slide.
        """
        try:
            from pptx import Presentation
            from pptx.util import Emu
            from PIL import ImageDraw, ImageFont
            import io

            prs = Presentation(str(pptx_path))
            slide_width = prs.slide_width or Emu(9144000)
            slide_height = prs.slide_height or Emu(6858000)
            aspect = int(slide_height) / int(slide_width)
            height = int(width * aspect)

            pages: List[Tuple[int, tuple]] = []

            for idx, slide in enumerate(prs.slides, start=1):
                text_parts: List[str] = []
                embedded_images: List[Image.Image] = []
                img_idx = 0
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                text_parts.append(line)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip(" |"):
                                text_parts.append(row_text)
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            blob = shape.image.blob
                            pil_img = Image.open(io.BytesIO(blob))
                            pil_img.load()
                            embedded_images.append(pil_img)
                            text_parts.append(f"![](images/{img_idx}.jpg)")
                            img_idx += 1
                        except Exception:
                            pass

                text = "\n".join(text_parts)

                img = Image.new("RGB", (width, height), color="white")
                draw = ImageDraw.Draw(img)
                try:
                    font = ImageFont.truetype("arial.ttf", 18)
                except Exception:
                    font = ImageFont.load_default()

                y_offset = 20
                for line in text_parts:
                    if y_offset > height - 30:
                        break
                    draw.text((20, y_offset), line[:120], fill="black", font=font)
                    y_offset += 24

                pages.append((idx, (text, img, embedded_images)))

            return pages

        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX parsing. "
                "Install with: pip install python-pptx"
            )

    @staticmethod
    def get_pptx_page_count(pptx_path: Path) -> int:
        """Get the number of slides in a PPTX file."""
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            return len(prs.slides)

        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX parsing. "
                "Install with: pip install python-pptx"
            )


def extract_document_pages(
    file_path: Path, prefer_text: bool = False
) -> List[Tuple[int, str | Image.Image]]:
    """Extract pages from document (PDF/EPUB/PPTX).

    Args:
        file_path: Path to document file
        prefer_text: If True, extract text directly when possible

    Returns:
        List of tuples (page_number, content) where content is text or Image
    """
    parser = DocumentParser()

    if parser.is_pdf(file_path):
        if prefer_text:
            try:
                return parser.parse_pdf_text_and_images(file_path)
            except Exception:
                return parser.parse_pdf_to_images(file_path)
        else:
            return parser.parse_pdf_to_images(file_path)

    elif parser.is_epub(file_path):
        if prefer_text:
            return parser.parse_epub_text_and_images(file_path)
        else:
            return parser.parse_epub_to_images(file_path)

    elif parser.is_pptx(file_path):
        if prefer_text:
            return parser.parse_pptx_text_and_images(file_path)
        else:
            return parser.parse_pptx_to_images(file_path)

    else:
        raise ValueError(f"Unsupported file format: {file_path.suffix}")


def iter_document_pages(
    file_path: Path,
    prefer_text: bool = False,
    start_page: int = 1,
) -> Iterator[Tuple[int, str | Image.Image | tuple[str, Image.Image]]]:
    """Yield document pages without materializing the whole document in memory."""
    parser = DocumentParser()

    if parser.is_pdf(file_path):
        if prefer_text:
            yield from parser.iter_pdf_text_and_images(file_path, start_page=start_page)
        else:
            yield from parser.iter_pdf_images(file_path, start_page=start_page)
        return

    pages = extract_document_pages(file_path, prefer_text=prefer_text)
    for page_num, content in pages:
        if page_num >= start_page:
            yield page_num, content


def get_document_page_count(file_path: Path) -> int:
    """Get the total number of logical pages/items in a document."""
    parser = DocumentParser()

    if parser.is_pdf(file_path):
        return parser.get_pdf_page_count(file_path)
    if parser.is_epub(file_path):
        return parser.get_epub_page_count(file_path)
    if parser.is_pptx(file_path):
        return parser.get_pptx_page_count(file_path)

    raise ValueError(f"Unsupported file format: {file_path.suffix}")
